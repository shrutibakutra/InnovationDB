import tldextract
import io
from random import randint, random
import os
import urllib.request
import re
import pandas as pd
import cv2
import requests
from striprtf.striprtf import rtf_to_text
from api.utils import remove_punc_specialChar
from textblob.classifiers import NaiveBayesClassifier
from api.utils import url_classification
import csv
BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
from rest_framework.views import APIView
from rest_framework.response import Response
import json
import PyPDF2
import os.path
import pandas
import docx2txt
import collections
import textract
from pptx import Presentation
from api.utils import Text_Comparision_controller
from pptx import Presentation as ABC
from pyexcel_ods import get_data
from odf import text, teletype
from odf.opendocument import load
from django.core.files.storage import FileSystemStorage
from pathlib import Path
import cv2 as cv
import numpy as np
import numpy
import os
from PIL.ExifTags import TAGS
import piexif
from nltk.corpus import stopwords
from nltk.stem import PorterStemmer
from nltk.tokenize import word_tokenize, sent_tokenize
from piexif import ExifIFD
from PIL import Image, ImageEnhance
from api.utils import webpage_to_text, Scrap_google_result, \
    find_cosine_similarity, get_companies_details_from_googlemaps, \
    webpage_to_text_no_clean, find_similarity_keyword, webpage_to_text_p, webpage_to_text_ul, \
    find_web_to_text_div_and_p, get_from_google_maps, webpage_to_text_li, find_keyword_repetition, \
    find_cosine_similarity_with_keyword_repetition, find_rating, webpage_to_text_find_div_and_P, \
    text_classification_by_repeated_keywords, find_classification, Scrap_web_result, find_html_to_text, find_uniquePhrases, \
    find_rating_percentage_list, find_weburl_classification_no_autoclassification, find_classification_no_auto, replica_similarity_keyword,replica_cosine_similarity,  \
    text_url_classification, url_to_text, Documen_job,Scrap_web_result_multisearch,find_sentence_similarity,Bing_scraper ,Google_scraper,find_word_similarity

from searchBar.views import delete_duplicate_link
from searchBar.models import KeywordSearch, Keyword_category , SearchResult
import wordninja
from .models import Links

class json_converter(APIView):

    def post(self, request):
        if 'filename' in request.data:
            inputfile = request.data['filename']
            filename = inputfile.name
            if filename.endswith('.txt'):
                glob = self.txt(inputfile)
            elif filename.endswith('.rtf'):
                glob = self.rtf(inputfile)
            elif filename.endswith('.pdf'):
                glob = self.pdf(inputfile)
            elif filename.endswith('.csv'):
                glob = self.csv(inputfile)
            elif filename.endswith('.odt'):
                glob = self.odt(inputfile)
            elif filename.endswith('.ods'):
                glob = self.ods(inputfile)
            elif filename.endswith('.odp'):
                glob = self.odp(inputfile)
            elif filename.endswith('.doc'):
                glob = self.doc(inputfile)
            elif filename.endswith('.docx'):
                glob = self.docx(inputfile)
            elif filename.endswith('.xls'):
                glob = self.xls(inputfile)
            elif filename.endswith('.xlsx'):
                glob = self.xlsx(inputfile)
            elif filename.endswith('.xlr'):
                glob = self.xlr(inputfile)
            elif filename.endswith('.pps'):
                glob = self.pps(inputfile)
            elif filename.endswith('.ppsx'):
                glob = self.ppsx(inputfile)
            elif filename.endswith('.ppt'):
                glob = self.ppt(inputfile)
            elif filename.endswith('.pptx'):
                glob = self.pptx(inputfile)
            return Response(glob)
        elif 'fileLink' in request.data:
            inputfile = request.data['fileLink']
            if inputfile.endswith('.txt'):
                data = urllib.request.urlopen(inputfile)
                text = ''
                for line in data:
                    decoded_line = line.decode("utf-8")
                    text += decoded_line
                data = self.textclean(text)
                return Response({'data': data})
            elif inputfile.endswith('.PDF') or inputfile.endswith('.pdf'):
                try:
                    r = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    f = io.BytesIO(r.content)
                    pdfFile = PyPDF2.PdfFileReader(f)
                    No_of_pages = pdfFile.getNumPages()
                    pages_content = ''
                    for i in range(No_of_pages):
                        page = pdfFile.getPage(i)
                        page_content = page.extractText()
                        pages_content += page_content
                    test = self.textclean(pages_content)
                    return Response({'data': test})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.csv'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    print(myfile)
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.txt')
                    open(file_path, 'wb').write(myfile.content)
                    source = open(file_path, "r")
                    decoded_file = source.read()
                    data = self.textclean(decoded_file)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.odt'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.odt')
                    open(file_path, 'wb').write(myfile.content)
                    text = textract.process(file_path)
                    esr = text.decode(encoding="utf-8")
                    data = self.textclean(esr)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.ods'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.ods')
                    open(file_path, 'wb').write(myfile.content)
                    data1 = get_data(file_path)
                    alphanumeric = [character for character in str(data1) if
                                    character.isalnum() or character.isspace()]
                    alphanumeric = "".join(alphanumeric)
                    data = self.textclean(alphanumeric)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.odp'):
                try:
                    from odf import text, teletype
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.odp')
                    open(file_path, 'wb').write(myfile.content)
                    textdoc = load(file_path)
                    allparas = textdoc.getElementsByType(text.P)
                    linelenght = len(allparas)
                    texts = ''
                    for line in range(linelenght):
                        test = teletype.extractText(allparas[line])
                        texts += test
                    data = self.textclean(texts)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.docx'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.docx')
                    open(file_path, 'wb').write(myfile.content)
                    text = docx2txt.process(file_path)
                    data = self.textclean(text)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.doc'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    data = urllib.request.urlopen(inputfile)
                    with open('media/temp/temp.doc', 'wb') as output:
                        output.write(data.read())
                    text = textract.process('media/temp/temp.doc')
                    text = text.decode("utf-8")
                    data = self.textclean(text)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.xlsx'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.xlsx')
                    open(file_path, 'wb').write(myfile.content)
                    df = pandas.read_excel(file_path)
                    df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
                    # data = self.textclean(text)
                    return Response({'data': df})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.xls'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.xls')
                    open(file_path, 'wb').write(myfile.content)
                    df = pandas.read_excel('media/temp/temp.xls')
                    df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
                    # data = self.textclean(df)
                    return Response({'data': df})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.pptx'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.pptx')
                    open(file_path, 'wb').write(myfile.content)
                    prs = Presentation(file_path)
                    text_runs = ''
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    text_runs += run.text
                    data = self.textclean(text_runs)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.xlr'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.xlr')
                    xlrfile = pandas.read_excel('media/temp/temp.xlr')
                    json_str = xlrfile.to_json(orient='records')
                    data = self.textclean(json_str)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
            elif inputfile.endswith('.rtf'):
                try:
                    myfile = requests.get(inputfile, headers={'User-Agent': 'Mozilla/5.0'})
                    file_path = os.path.join(BASE_DIR, 'media/temp/temp.rtf')
                    open(file_path, 'wb').write(myfile.content)
                    text = open(file_path, 'r')
                    rtf = text.read()
                    txt = rtf_to_text(rtf)
                    data = self.textclean(txt)
                    return Response({'data': data})
                except:
                    return Response({'data': ' '})
        else:
            return Response({'massage': 'Please pass file or Link'})

    def textclean(self, text):
        remove_n = text.replace('\n', ' ')
        remove_r = remove_n.replace('\r', ' ')
        remove_t = remove_r.replace('\t', ' ')
        clean_text = re.sub(' +', ' ', remove_t)
        return remove_t

    def txt(self, txt):
        filen = str(txt)
        line = txt.readlines()
        txtdata = {}
        text = ''
        for Line in line:
            decode_text = (Line.decode("utf-8"))
            text += decode_text
        txt = json.dumps(txtdata)
        formated_text = remove_punc_specialChar(text)
        save_file = open("media/jsonresponse/" + filen, "w")
        json.dump(txt, save_file, indent=4, sort_keys=False)
        save_file.close()
        data = {"data": formated_text}
        return data

    def pdf(self, pdf):
        filen = str(pdf)
        textfile = Path(filen).stem + ".txt"
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/' + textfile, pdf)
        test = open("media/jsonresponse/" + textfile, "w")
        read_pdf = PyPDF2.PdfFileReader(pdf)
        number_of_pages = read_pdf.getNumPages()
        text = ''
        for page_number in range(number_of_pages):  # use xrange in Py2
            page = read_pdf.getPage(page_number)
            page_content = page.extractText()
            test.write(page_content)
            text += page_content
        # formated_text = remove_punc_specialChar(text)
        # data = {"data": text}
        # print(type(data))
        return text

    def rtf(self, rtf):
        filen = str(rtf)
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/' + filen, rtf)
        textfile = Path(filen).stem + ".txt"
        save_file = open("media/jsonresponse/" + textfile, "w")
        path = "media/jsonresponse/" + filen
        test = ''
        with open(path, 'r') as file:
            rtf = file.read()
            txt = rtf_to_text(rtf)
            save_file.write(txt)
            test += txt
        print(txt)
        os.remove("media/jsonresponse/" + filen)
        formated_text = remove_punc_specialChar(test)
        data = {"data": formated_text}
        return data

    def csv(self, csv):
        filen = str(csv)
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/' + filen, csv)
        textfile = Path(filen).stem + ".txt"
        save_file = open("media/jsonresponse/" + textfile, "w")
        path = "media/jsonresponse/" + filen
        test = ''
        with open(path, 'r') as file:
            text = file.read()
            csv = rtf_to_text(text)
            save_file.write(csv)
            test += csv
        os.remove("media/jsonresponse/" + filen)
        formated_text = remove_punc_specialChar(test)

        data = {"data": formated_text}
        return data

    def xlsx(self, xlsx):
        filen = str(xlsx)
        xlsxfile = pd.read_excel(xlsx)
        return xlsxfile.to_string()
        # json_str = xlsxfile.to_json(orient='records')
        # textfile = Path(filen).stem + ".txt"
        # save_file = open("media/jsonresponse/" + textfile, "w")
        # json.dump(json_str, save_file, indent=4, sort_keys=False)
        # save_file.close()
        # formated_text = remove_punc_specialChar(json_str)
        # print(formated_text)
        # data = str(formated_text)
        # return json_str

    def xls(self, xls):
        filen = str(xls)
        xlsfile = pandas.read_excel(xls)
        json_str = xlsfile.to_json(orient='records')
        textfile = Path(filen).stem + ".txt"
        save_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(json_str, save_file, indent=4, sort_keys=False)
        save_file.close()
        formated_text = remove_punc_specialChar(json_str)

        data = {"data": formated_text}
        return data

    def xlr(self, xlr):
        filen = str(xlr)
        xlrfile = pandas.read_excel(xlr)
        json_str = xlrfile.to_json(orient='records')
        textfile = Path(filen).stem + ".txt"
        save_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(json_str, save_file, indent=4, sort_keys=False)
        save_file.close()
        formated_text = remove_punc_specialChar(json_str)

        data = {"data": formated_text}
        return data

    def docx(self, docx):
        filen = str(docx)
        text = docx2txt.process(docx)
        li = [x for x in text.split('\n')]
        li1 = list(filter(None, li))
        li = ''.join(li1)
        json_li = []
        for x in li:
            x = x[2:]
            y = x.split(',')
            d = collections.defaultdict()
            for m in y:
                z = m.split(':')
                z1 = [x.strip() for x in z]
                d[z1[0]] = z1[0]
            json_li.append(d)
        textfile = Path(filen).stem + ".txt"
        out_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(li, out_file, indent=4, sort_keys=False)
        print(li)
        out_file.close()
        formated_text = remove_punc_specialChar(li)
        data = {"data": formated_text}
        return data

    def doc(self, doc):
        filen = str(doc)
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/' + filen, doc)
        textfile = Path(filen).stem + ".txt"
        import textract
        save_file = open("media/jsonresponse/" + textfile, "w")
        text = textract.process("media/jsonresponse/" + filen)
        text = text.decode("utf-8")
        save_file.write(text)
        os.remove("media/jsonresponse/" + filen)
        formated_text = remove_punc_specialChar(text)
        data = {"data": formated_text}
        return data

    def ods(self, ods):
        filen = str(ods)
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/' + filen, ods)
        data1 = get_data("media/jsonresponse/" + filen)
        textfile = Path(filen).stem + ".txt"
        out_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(data1, out_file, indent=4, sort_keys=False)
        os.remove("media/jsonresponse/" + filen)
        out_file.close()
        alphanumeric = [character for character in str(data1) if
                        character.isalnum() or character.isspace()]

        alphanumeric = "".join(alphanumeric)
        data = {"data": alphanumeric}
        return data

    def odt(self, odt):
        filen = str(odt)
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/*.odt', odt)
        text = textract.process("media/jsonresponse/*.odt")
        esr = text.decode(encoding="utf-8")
        textfile = Path(filen).stem + ".txt"
        out_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(esr, out_file, indent=4, sort_keys=False)
        out_file.close()
        os.remove("media/jsonresponse/*.odt")
        formated_text = remove_punc_specialChar(esr)
        data = {"data": formated_text}
        return data

    def odp(self, odp):
        filen = str(odp)
        textdoc = load(odp)
        allparas = textdoc.getElementsByType(text.P)
        linelenght = len(allparas)
        textfile = Path(filen).stem + ".txt"
        f = open("media/jsonresponse/" + textfile, "w")
        texts = ''
        for line in range(linelenght):
            test = teletype.extractText(allparas[line])
            texts += test
            json.dump(test, f, indent=4, sort_keys=False)
        data = {'data': texts}
        return data

    def pptx(self, pptx):
        filen = str(pptx)
        test = pptx
        prs = Presentation(test)
        prs.save('media/jsonresponse/' + filen)
        prs = Presentation("media/jsonresponse/" + filen)
        text_runs = ''
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs += run.text
        textfile = Path(filen).stem + ".txt"
        out_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(text_runs, out_file, indent=4, sort_keys=False)
        os.remove("media/jsonresponse/" + filen)
        formated_text = remove_punc_specialChar(text_runs)

        data = {"data": text_runs}
        return data

    def ppt(self, ppt):
        filen = str(ppt)
        fs = FileSystemStorage()
        filename = fs.save('jsonresponse/*.pptx', ppt)
        prs = ABC("media/jsonresponse/*.pptx")
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        textfile = Path(filen).stem + ".txt"
        out_file = open("media/jsonresponse/" + textfile, "w")
        json.dump(text_runs, out_file, indent=4, sort_keys=False)
        os.remove("media/jsonresponse/*.pptx")
        formated_text = remove_punc_specialChar(text_runs)

        data = {"data": formated_text}
        return data

    def pps(self, pps):
        try:
            filen = str(pps)
            fs = FileSystemStorage()
            filename = fs.save('jsonresponse/*.pptx', pps)
            prs = ABC("media/jsonresponse/*.pptx")
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            textfile = Path(filen).stem + ".txt"
            out_file = open("media/jsonresponse/" + textfile, "w")
            json.dump(text_runs, out_file, indent=4, sort_keys=False)
            os.remove("media/jsonresponse/*.pptx")
            formated_text = remove_punc_specialChar(text_runs)

            data = {"data": formated_text}
            return data
        except:
            data = {"error": "Your file  is not a PowerPoint file"}
            return data

    def ppsx(self, ppsx):
        try:
            filen = str(ppsx)
            test = ppsx
            prs = Presentation(test)
            prs.save('media/jsonresponse/' + filen)
            prs = Presentation("media/jsonresponse/" + filen)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            text_runs.append(run.text)
            textfile = Path(filen).stem + ".txt"
            out_file = open("media/jsonresponse/" + textfile, "w")
            json.dump(text_runs, out_file, indent=4, sort_keys=False)
            os.remove("media/jsonresponse/" + filen)
            formated_text = remove_punc_specialChar(text_runs)
            data = {"data": formated_text}
            return data
        except:
            data = {"error": "Your file  is not a PowerPoint file"}
            return data


class text_to_converter(APIView):
    def post(self, request):
        all_text = request.data
        result = Text_Comparision_controller(all_text)
        return Response(result)


class similarity_keyword(APIView):
    def post(self, request):
        percentage_value = find_similarity_keyword(request.data['text'], request.data['text2'])
        return Response(percentage_value)


class Similarity_cosine(APIView):
    def post(self, request):
        similarity_result = find_cosine_similarity(request.data['text1'], request.data['text2'])
        return Response({'data': similarity_result})


class Image_manuplator(APIView):
    def post(self, request):
        try:
            image = request.data['image']
            alpha_val = request.data['alpha']
            beta_val = request.data['beta']
        except:
            return Response({'message': 'Value is missing'})
        ext = image.name.split(".")[-1]
        random_digit = randint(1000, 9999)
        image_renamed = str(random_digit) + '.' + ext
        fs = FileSystemStorage()
        fs.save('temp_images/' + image_renamed, image)
        file_path = os.path.join(BASE_DIR, 'media/temp_images/' + image_renamed)
        storing_path = os.path.join(BASE_DIR, 'media/manuplated_images/' + image_renamed)

        # Read image given by user
        image = cv.imread(file_path)
        if image is None:
            return Response({'massage': 'Could not open or find the image'})

        new_image = np.zeros(image.shape, image.dtype)
        alpha = 1.0  # Simple contrast control
        beta = 0  # Simple brightness control
        # Initialize values
        try:
            alpha = float(alpha_val)  # float(input('* Enter the alpha value [1.0-3.0]: '))
            beta = int(beta_val)  # int(input('* Enter the beta value [0-100]: '))
        except ValueError:
            return Response({'message': 'Error, not a number'})
        for y in range(image.shape[0]):
            for x in range(image.shape[1]):
                for c in range(image.shape[2]):
                    new_image
                    [y, x, c] = np.clip(alpha * image[y, x, c] + beta, 0, 255)
        cv2.imwrite(storing_path, new_image)
        os.remove(file_path)
        return Response({'massage': 'Image Saved Successfully'})


class Image_change(APIView):
    def post(self, request):
        try:
            uploadfile = request.data['image']
        except:
            return Response({'message': 'Value is missing'})
        image_file = "media/sample_images/sample.jpg"
        image = Image.open(image_file)
        np_im = numpy.array(image)
        inth = np_im.shape[0]
        intw = np_im.shape[1]
        exif = {}
        for tag, value in image._getexif().items():
            if tag in TAGS:
                exif[TAGS[tag]] = value
        ExifResolutionUnit = exif['ResolutionUnit']
        print(ExifResolutionUnit)
        XResolution = exif['XResolution']
        print(XResolution)
        YResolution = exif['YResolution']
        print(YResolution)
        Sharpness = exif['Sharpness']
        print(Sharpness)
        ExposureBiasValue = exif['ExposureBiasValue']
        print(ExposureBiasValue)
        int1 = XResolution[0]
        int2 = XResolution[1]
        int3 = YResolution[0]
        int4 = YResolution[1]
        zeroth_ifd = {
            piexif.ImageIFD.XResolution: (int1, int2),
            piexif.ImageIFD.YResolution: (int3, int4)
        }
        exif_ifd = {
            piexif.ExifIFD.Sharpness: 63455,
            ExifIFD.ExposureBiasValue: ExposureBiasValue,
        }
        first_ifd = {
            piexif.ImageIFD.XResolution: (int1, int2),
            piexif.ImageIFD.YResolution: (int3, int4),
        }
        exif_dict = {"0th": zeroth_ifd, "Exif": exif_ifd}
        exif_bytes = piexif.dump(exif_dict)
        ext = uploadfile.name.split(".")[-1]
        test = str(uploadfile)
        random_digit = randint(1000, 9999)
        image_renamed = str(random_digit) + '.' + ext
        im = Image.open(uploadfile)
        im.save("media/manuplated_images/" + "change" + test, exif=exif_bytes)
        im1 = Image.open("media/manuplated_images/" + "change" + test)
        width = intw
        height = inth
        imge = im1.resize((width, height), Image.NEAREST)
        imge.save("media/manuplated_images/" + "change" + test)
        imagec = Image.open("media/manuplated_images/" + "change" + test)
        converter = ImageEnhance.Color(imagec)
        img = converter.enhance(ExifResolutionUnit)
        # imgBrightness = ImageEnhance.Brightness(imagec)
        # img = imgBrightness.enhance(1)
        imgContrast = ImageEnhance.Contrast(imagec)
        img = imgContrast.enhance(ExifResolutionUnit)
        imgSharpness = ImageEnhance.Sharpness(img)
        img = imgSharpness.enhance(ExifResolutionUnit)

        test = img.save("media/manuplated_images/" + "change" + test)
        return Response({'massage': 'Image Saved Successfully'})


class text_to_paraphrase(APIView):
    def _create_frequency_table(self, text_string) -> dict:
        stopWords = set(stopwords.words("english"))
        words = word_tokenize(text_string)
        ps = PorterStemmer()

        freqTable = dict()
        for word in words:
            word = ps.stem(word)
            if word in stopWords:
                continue
            if word in freqTable:
                freqTable[word] += 1
            else:
                freqTable[word] = 1

        return freqTable

    def _score_sentences(self, sentences, freqTable) -> dict:
        sentenceValue = dict()

        for sentence in sentences:
            word_count_in_sentence = (len(word_tokenize(sentence)))
            for wordValue in freqTable:
                if wordValue in sentence.lower():
                    if sentence[:10] in sentenceValue:
                        sentenceValue[sentence[:10]] += freqTable[wordValue]
                    else:
                        sentenceValue[sentence[:10]] = freqTable[wordValue]

            sentenceValue[sentence[:10]] = sentenceValue[sentence[:10]] // word_count_in_sentence

        return sentenceValue

    def _find_average_score(self, sentenceValue) -> int:
        sumValues = 0
        for entry in sentenceValue:
            sumValues += sentenceValue[entry]

        # Average value of a sentence from original text
        average = int(sumValues / len(sentenceValue))

        return average

    def _generate_summary(self, sentences, sentenceValue, threshold):
        sentence_count = 0
        summary = ''

        for sentence in sentences:
            if sentence[:10] in sentenceValue and sentenceValue[sentence[:10]] > (threshold):
                summary += " " + sentence
                sentence_count += 1

        return summary

    def post(self, request):
        text = request.data['text']
        # 1 Create the word frequency table
        freq_table = self._create_frequency_table(text)

        # 2 Tokenize the sentences
        sentences = sent_tokenize(text)

        # 3 Important Algorithm: score the sentences
        sentence_scores = self._score_sentences(sentences, freq_table)

        # 4 Find the threshold
        threshold = self._find_average_score(sentence_scores)

        # 5 Important Algorithm: Generate the summary
        summary = self._generate_summary(sentences, sentence_scores, 1.5 * threshold)

        return Response({'paraphrase_text': summary})


class text_classification(APIView):
    def post(self, request):
        text = request.data['text']

        train = [
            ('Technology ,latest technology,tech,science technology ', 'TECHNOLOGY'),
            ('Business , market , business plan ', 'Business'),
            ('Music ,music ,speaker ,music sound ', 'Audio'),
            ('wearable, hearables, astronautics ,technologist ,cyberphobia , animatronics ', 'Wearable Technology'),
            ('VR , artificial reality , virtual , virtuality , virtually ', 'Virtual and Augmented Reality'),
            ('videoshoot, camera , video shooting , video shoot , videography ', 'Videography'),
            ('telecom , Internet , telecommunication ', 'Telecommunications and Networks'),
            ('game,play, player , sport , sporting , sports , turf', 'Sport'),
            ('computer software , application software , install,  developer , uninstall ,search engine ', 'Software'),
            ('security , safety , safe , secure ,security guard ', 'Safety'),
            ('roboteer , animatronics , cloud robotics , robot', 'Robotics'),
            ('clinic , patient , medical, medicine, surgeon , dispensary , medically , heart, donor, transplants, '
             'hospital, cardiac, lungs,'
             'first aid, cardiac pump, artificial heart ', 'Medical'),
            ('premarketing ,viral marketing , telemarketing , network marketing ,advertising , social marketing  ',
             'Marketing'),
            ('industrial , industrialize , Manufacturing ,  Manufactur ', 'Manufacturing and Industry'),
            ('graphic , graphically , diagram , design ', 'Graphic'),
            ('financial , financier , liability ', 'Financial'),
            ('environmentalist , nature , tree ,forest , environmentalism , ', 'Environment'),
            ('kinetic energy , atomic energy ,solar energy , energy ', 'Energy'),
            ('electrical engineer ,electrotechnology, electrician , electrical , electricity ',
             'Electrical Engineering'),
            ('transportation , transport, transportable, transportation vehicles, vehicles, taxi, bus, jet pack, '
             'rocket, rocket pack, jet, rocket belt ', 'Transport'),
            ('school ,educative ,institute , study , college ', 'Educational'),
            ('home electronics ,home electronics , electronic ', 'Consumer Electronics'),
            ('constructive ,structure ,build , architect ', 'Construction'),
            ('computation , technological , hackerspace ', 'Computing and Technology'),
            ('botanical , botany , farming ', 'Botany'),
        ]

        cl = NaiveBayesClassifier(train)

        alphanumeric = [character for character in text if
                        character.isalnum() or character.isspace()]

        alphanumeric = "".join(alphanumeric)
        string = alphanumeric.replace('\r', '').replace('\n', ' ')
        text_tokens = word_tokenize(string)
        tokens_without_sw = [word for word in text_tokens if not word in stopwords.words()]
        clean_text = ''
        for word in tokens_without_sw:
            clean_text += ' ' + word
        test = cl.classify(text)

        return Response({'paraphrase_text': test})


class clean_text(APIView):
    def post(self, requests):
        text = requests.data['text']
        clean_texts = remove_punc_specialChar(text)
        return Response({'data': clean_texts})


class document_classification(APIView):
    def post(self, request):
        text = request.data['text']
        train = [
            (
                'Technical,Technology ,latest technology,tech,science  technology , technological ,mechanical , '
                'engineering, specialized ,scientific , practical , operational , expert , specialist , tactical , '
                'engineer , training ,performance ,physical , analytical , fundamental , specific , proficient ',
                'TECHNICAL'),
            (
                'investment,bank,economics,stock, capital,banking,finances,economic,management,finance,accounting,'
                'economy,funds,corporate,equity,insurance,business,financial engineering,financial economics,'
                'financial market,asset,financing,interest,credit,banks,budget,money,investment management,'
                'personal finance',
                'FINANCIAL'),
            (
                'advertising ,mercantile,private,commercialized,commercial message,mercenary,moneymaking,inferior,'
                'advertisement,commerce,infomercial,market,operators,marketing,trade,industrial,retail,corporate,sales,'
                'industry,commercialism ,export ,product,shopping,industry,merchant,retail,economic,shopping',
                'COMMERCIAL'),
            (
                'research,scientist,scientific method,physics,mathematics,knowledge,astronomy,medicine,chemistry,'
                'biology, '
                'experiment,logic,observation,theory,academic,mathematical,natural science,evolution,academy,'
                'scientists, '
                'scientific,Experperiment,laboratory,biochemist ,physics , physician , mineralogist ',
                'SCIENTIFIC'),
        ]
        cl = NaiveBayesClassifier(train)
        test = cl.classify(text)
        print(test)

        return Response({'Document-type': test})


class web_text_classification(APIView):
    def post(self, request):
        classified_text = text_classification_by_repeated_keywords(request.data)  # find_web_text_classification
        return Response({'Document-type': classified_text})


class Website_scraping(APIView):
    def post(self, request):
        keyword = request.data['keyword']
        num_pages = request.data['num_pages']
        data = Scrap_google_result(keyword, num_pages)
        return Response({'links': data})


class website_to_text(APIView):
    def post(self, request):
        url = request.data['url']
        text = webpage_to_text(url)
        return Response({'data': text})


"""Website To Text New Flow SOme Temp Api's"""


class website_to_text_no_clean(APIView):
    def post(self, request):
        url = request.data['url']
        text = webpage_to_text_no_clean(url)
        return Response({'data': text})


class website_to_text_p(APIView):
    def post(self, request):
        url = request.data['url']
        text = webpage_to_text_p(url)
        return Response({'data': text})


class website_to_text_ul(APIView):
    def post(self, request):
        url = request.data['url']
        text = webpage_to_text_ul(url)
        return Response({'data': text})


class website_to_text_li(APIView):
    def post(self, request):
        url = request.data['url']
        text = webpage_to_text_li(url)
        return Response({'data': text})


"""End Of Flow"""


class delete_duplicatelink(APIView):
    def post(self, request):
        id = request.data['keywordId_id']
        print(id)
        deleteurl = delete_duplicate_link(id)
        return Response({'data': 'sucess'})


class url_to_domain_name(APIView):
    def post(self, request):
        url = request.data['url']
        info = tldextract.extract(url)
        """Full Information For any url"""
        # domainInfo = {'domain_name' : info.domain,
        #                 "subdomain_name" : info.subdomain,
        #                 "suffix": info.suffix,
        #                 "domain_with_suffix": info.registered_domain,
        #                 "full_domain" : '.'.join(info)
        #               }
        return Response({'domain_name': info.domain})


class weburl_classification(APIView):
    def post(self, request):
        url = request.data['url']
        text = url_classification(url)
        return Response({'data': text})


class weburl_classification_no_autoclassification(APIView):
    def post(self, request):
        url = request.data['url']
        text = find_weburl_classification_no_autoclassification(url)
        return Response({'data': text})


class weburl_classification_list(APIView):
    def post(self, request):
        urls = request.data['url']
        data = []
        for url in urls:
            text = url_classification(url)
            texts = {url: text}
            data.append(texts)
        return Response({'data': data})


class splittext_withoutspaces(APIView):
    def post(self, request):
        text = request.data['text']
        split = wordninja.split(text)
        splittext = ' '.join(map(str, split))
        return Response({'splittext ': splittext})


class identify_companies_from_googlemaps(APIView):
    def post(self, request):
        company_details = get_companies_details_from_googlemaps(request.data)

        return Response({'company_info ': company_details})


class google_maps(APIView):
    def post(self, request):
        company_details = get_from_google_maps(request.data)

        return Response({'company_info ': company_details})


class find_words(APIView):
    def post(self, request):
        list_of_words = []
        if request.data['words']:
            words = request.data['words']
            words = re.sub('[^a-zA-Z0-9\n\-\.]', ',', words)
            list_of_words = words.split(',')
        # print(list_of_words)
        else:
            file = request.data['file']
            line = file.readlines()
            for Line in line:
                for word in Line.split():
                    decode_text = (word.decode("utf-8"))
                    list_of_words.append(decode_text)

        list_of_text = []
        text = request.data['conntent']
        for word in text.split():
            list_of_text.append(word)

        comman = set(list_of_words) & set(list_of_text)
        similar_word = set(sorted(comman, key=lambda k: list_of_words.index(k)))

        if similar_word:
            return Response({'words': similar_word})
        else:
            return Response({'words': 'no similar words found'})


class url_to_domain_suffix(APIView):
    def post(self, request):
        url = request.data['url']
        info = tldextract.extract(url)
        domain_name = '.'.join(info)
        return Response({'domain': domain_name})


class keyword_repetition(APIView):
    def post(self, request):
        similarity = find_keyword_repetition(request.data['text1'], request.data['text2'])
        return Response(similarity)


class cosine_similarity_with_keyword_repetition(APIView):
    def post(self, request):
        similarity = find_cosine_similarity_with_keyword_repetition(request.data['text1'], request.data['text2'])
        return Response(similarity)


class web_to_text_div_and_p(APIView):
    def post(self, request):
        text = find_web_to_text_div_and_p(request.data['url'])
        return Response(text)


class rating(APIView):
    def post(self, request):
        rating_result = find_rating(request.data['cosine_value'], request.data['keyword_repetition_value'])
        return Response(rating_result)

class classification(APIView):
    def post(self, request):
        classification_result = find_classification(request.data['urls'])
        return Response({'classification_result': classification_result})


# Testing purpouse Demo replica
class classification_no_auto(APIView):
    def post(self, request):
        classification_result = find_classification_no_auto(request.data['urls'])
        return Response({'classification_result': classification_result})

class Bing_Website_scraper(APIView):
    def post(self, request):
        scrap_result = Bing_scraper(request.data['keyword'], request.data['num_pages'])
        return Response({'result': scrap_result})
        
class Website_scraper(APIView):
    def post(self, request):
        scrap_result = Scrap_web_result(request.data['keyword'], request.data['num_pages'])
        return Response({'result': scrap_result})


class Multi_search(APIView):
     def post(self, request):
        scrap_result = Scrap_web_result_multisearch(request.data['keyword'], request.data['num_pages'])
        return Response({'result': scrap_result})

class delete_duplicate_dict(APIView):
    def post(self, request):
        remove_duplicate_data = pd.DataFrame(request.data['list']).drop_duplicates("link").to_dict('records')
        return Response({'removed_duplicate_data': remove_duplicate_data})


class html_to_text(APIView):
    def post(self, request):
        html_markdown_text = find_html_to_text(request.data['url'])
        return Response({'text': html_markdown_text})


class rating_percentage_list(APIView):
    def post(self, request):
        ratings = find_rating_percentage_list(request.data['urls'], request.data['referance_text'])
        return Response({'rating': ratings})


class text_classification_url_and_text(APIView):
    def post(self, request):
        category = ''
        if request.data['text']:
            text = request.data['text']
            text = text_url_classification(text)
            category += text
        elif request.data['url']:
            url = request.data['url']
            text = find_html_to_text(url)
            text = text_url_classification(text)
            category += text
        return Response({'category': category})


class url_to_text_extractor(APIView):
    def post(self, request):
        text = url_to_text(request.data['url'])
        return Response({'text': text})


class auto_search_keyword(APIView):
    def post(self, request):
        if 'category' in request.data:
            categories = request.data['category']
            Key_word = request.data['title']
            Description = request.data['ref_text']
            types = request.data['type']
            categories.append('OTHERS')
        elif 'jsonfile' in request.data:
            file = request.data['jsonfile']
            filen = str(file)
            fs = FileSystemStorage()
            fs.save('media/temp/' + filen, file)
            f = open('media/media/temp/' + filen)
            data = json.load(f)
            Key_word = data["title"]
            categories = data["category"]
            Description = data["ref_text"]
            types = data["type"]
            f.close()
            os.remove('media/media/temp/' + filen)

        for type in types:
            store_Keyword = KeywordSearch(keyWord=Key_word, description=Description, filter=type, status=4)
            store_Keyword.save()
            for category in categories:
                store_category = Keyword_category(name=category.upper(), keywordId=store_Keyword)
                store_category.save()
        """Run BackGround Job For Panding Task"""
        Keyword_result = KeywordSearch.objects.filter(status=4).values()
        for filter in Keyword_result:
            del filter["created_at"]
            del filter["updated_at"]
            keyword_searchId = KeywordSearch.objects.get(pk=filter['id'])
            keyword_searchId.status = 0
            keyword_searchId.save()
            create_lists = [filter]
            Documen_job(create_lists)
        return Response({'Massage': 'keyword stored for crawling process'})


class Mutli_search(APIView):
    def post(self, request):
        keywords = request.data['keywords']
        categories = []
        types = []
        Description = ''
        num_pages = "20"
        links = []
        for keyword in keywords:
            scrap_link = Scrap_web_result(keyword,num_pages)
            result = {keyword : scrap_link}
            print(result)
            links.append(result)
        
        with open('media/Dataset/DEMO_Searches_-_Searches-1.csv', newline='') as csvfile:
            data = csv.DictReader(csvfile)
            for row in data:
                cat = row['CATEGORY'], row['SUBCATEGORY']
                types = row['TYPE']
                keywords = row['KEY 1'],row['KEY 2'], row['KEY 3'], row['KEY 4'] , row['KEY 5']
                for keyword in keywords:   
                    store_Keyword = KeywordSearch(keyWord=keyword,description=Description,filter=types, status=4)
                    store_Keyword.save()
                    for category in cat:
                        store_category = Keyword_category(name=category, keywordId=store_Keyword)
                        store_category.save()    

        Keyword_result = KeywordSearch.objects.filter(status=4).values()
        for filter in Keyword_result:
            del filter["created_at"]
            del filter["updated_at"]
            keyword_searchId = KeywordSearch.objects.get(pk=filter['id'])
            keyword_searchId.status = 0
            keyword_searchId.save()
            create_lists = [filter]
            print(create_lists)
            Documen_job(create_lists)

        return Response({'result': links})
class Similarity_word_match(APIView):
    def post(self, request):
        word_result = find_word_similarity(request.data['text1'], request.data['text2'])
        return Response({'data': word_result}
                )
class Similarity_sentence_match(APIView):
    def post(self, request):
        word_result = find_sentence_similarity(request.data['text1'], request.data['text2'])
        return Response({'data': word_result})

class Google_Website_scraper(APIView):
    def post(self, request):
        scrap_result = Google_scraper(request.data['keyword'], request.data['num_pages'])
        return Response({'result': scrap_result})


class Purpose_module(APIView):
    def post(self, request):
        text1 = request.data['Purpose Text']
        text2 = request.data['Comparison Text']
        word_result =replica_similarity_keyword(text1,text2)
        print(word_result)
        words = list(word_result)       
        lines = find_uniquePhrases(words,text2)
        test = replica_cosine_similarity(text1,lines)
        percentage = test['percentage']
        return Response({'percentage': percentage})

class Purpose_details(APIView):
    def post(self, request):
        text1 = request.data['Purpose Text']
        text2 = request.data['Comparison Text']
        word_result =replica_similarity_keyword(text1,text2)
        words = list(word_result)       
        lines = find_uniquePhrases(words,text2)
        test = replica_cosine_similarity(text1,lines)
        return Response({'result': test})

from selenium import webdriver    
from xvfbwrapper import Xvfb
import requests

class html_to_image(APIView):
    def post(self, request):
        urls = request.data['urls']
        result = []
        for url in urls:
            info = tldextract.extract(url)
            image_name = info.domain
            image_renamed = image_name + '.' + "png"
            destination="media/temp/"+image_renamed
            BASE = 'https://render-tron.appspot.com/screenshot/'
            path = destination
            ted = {url:"http://127.0.0.1:8000/" + path}
            result.append(ted)
            response = requests.get(BASE + url, stream=True)
            if response.status_code == 200:
                with open(path, 'wb') as file:
                    for chunk in response:
                        file.write(chunk)
        return Response({'image':result})


from searchBar.models import SearchResult

class SearchResults(APIView):
    def get(self,request):
        Results = SearchResult.objects.all().values()
        return Response(Results)

class Rating_Purpose(APIView):
    def post(self, request):
        urls = request.data['urls']
        Purpose = request.data['Purpose Text']
        referance_text = request.data['referance_text']
        ratings = find_rating_percentage_list( urls, referance_text)
        Purpose_rating = []
        
        for url in urls:
            text = find_html_to_text(url)
            word_result =replica_similarity_keyword(Purpose,text)
            words = list(word_result)       
            lines = find_uniquePhrases(words,text)
            cosine_similarity = replica_cosine_similarity(Purpose,lines)
            percentage = cosine_similarity['percentage']
            rat = {"url":url , "Purpose":percentage }
            Purpose_rating.append(rat)
        return Response({'rating': ratings , "Purpose_+_rating" : Purpose_rating })


class JsonSearch(APIView):
    def get(self,request):
        Results = KeywordSearch.objects.all().values()
        
        serializers = KeywordSearchSerializers(Results, many=True)
        return Response(serializers.data)


class Remove_links_table(APIView):
    def post(self,request):
        urls = request.data['links']
        Link = Links.objects.all().values()
        links = list(Link)
        database_link = []
        for link in links:
            url = link["link"]
            info = tldextract.extract(url)
            domain_name = info.domain
            database_link.append(domain_name)
        difference_link = []
        for url in urls:
            ext_domain = tldextract.extract(url)
            domain = ext_domain.domain
            if domain not in database_link:   
                difference_link.append(url)
        return Response({'Different_link':difference_link})


class Remove_links_table(APIView):
    def post(self,request):
        urls = request.data['links']
        Link = Links.objects.all().values()
        links = list(Link)
        database_link = []
        for link in links:
            url = link["link"]
            info = tldextract.extract(url)
            domain_name = info.domain
            database_link.append(domain_name)
        difference_link = []
        for url in urls:
            ext_domain = tldextract.extract(url)
            domain = ext_domain.domain
            if domain not in database_link:   
                difference_link.append(url)
        return Response({'Different_link':difference_link})
