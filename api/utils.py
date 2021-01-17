import difflib
import math
import operator
from collections import Counter
from string import digits
from textblob.classifiers import NaiveBayesClassifier
from django.db.models import Q
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
import requests
from bs4 import BeautifulSoup
import re
from sklearn.feature_extraction.text import TfidfTransformer
from sklearn.naive_bayes import MultinomialNB
from sklearn.feature_extraction.text import CountVectorizer
from sklearn.pipeline import Pipeline
from sklearn.model_selection import RandomizedSearchCV
from tldextract import tldextract
from nltk.stem import WordNetLemmatizer
import concurrent.futures
import logging
import csv
from id_interface import settings
from inscriptis import get_text

wnl = WordNetLemmatizer()
from searchBar.models import KeywordSearch, SearchResult, Keyword_category, company_info
import pandas as pd
import pathlib
import textract
import urllib.request
from odf.opendocument import load
from pyexcel_ods import get_data
import docx2txt
import pandas
import numpy as np
from striprtf.striprtf import rtf_to_text
from pptx import Presentation
import os
import io
import PyPDF2
from urllib.parse import urlparse

BASE_DIR = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))


def remove_punc_specialChar(text):
    remove_digits = str.maketrans('', '', digits)
    digit_removed = text.translate(remove_digits)
    alphanumeric = [character for character in digit_removed if
                    character.isalnum() or character.isspace() or character.isalpha()]

    alphanumeric = "".join(alphanumeric)
    string = alphanumeric.replace('\r', '').replace('\n', ' ')
    text_tokens = word_tokenize(string)

    tokens_without_sw = [word for word in text_tokens if not word in stopwords.words()]

    text = ''
    for word in tokens_without_sw:
        text += ' ' + word
    return text


def webpage_to_text(url):
    url = url
    res = requests.get(url)
    html_page = res.content
    soup = BeautifulSoup(html_page, 'html.parser')
    result = soup.find_all("div")
    blacklist = [
        '[document]',
        'noscript',
        'header',
        'html',
        'meta',
        'head',
        'input',
        'script',
        # there may be more elements you don't want, such as "style", etc.
    ]
    text_result = ''
    """Remove BlackListed Tag And extract Text"""
    for t in result:
        if t.parent.name not in blacklist:
            text = str(t.text)
            """Process For removing extra space The Text"""
            remove_n = text.replace('\n', ' ')
            clean_text = re.sub(' +', ' ', remove_n)
            remove_leading_space = clean_text.lstrip()
            text_result += remove_leading_space + ' '
    return text_result


"""Website To Text New Flow SOme Temp Api's"""


def webpage_to_text_no_clean(url):
    url = url
    res = requests.get(url)
    html_page = res.content
    soup = BeautifulSoup(html_page, 'html.parser')
    result = soup.find_all("div")
    blacklist = [
        '[document]',
        'noscript',
        'header',
        'html',
        'meta',
        'head',
        'input',
        'script',
        # there may be more elements you don't want, such as "style", etc.
    ]
    resultss = []
    """Remove BlackListed Tag And extract Text"""
    for t in result:
        if t.parent.name not in blacklist:
            text = str(
                t.text)
            """Process For Cleaning The Text"""
            remove_n = text.replace('\n', ' ')
            clean_text = re.sub(' +', ' ', remove_n)
            remove_leading_space = clean_text.lstrip()
            """SEt Tag name Of particular Text"""
            if not {t.parent.name: remove_leading_space} in resultss:
                if remove_leading_space and t.parent.name == 'div':
                    resultss.append({t.parent.name: remove_leading_space})
    return resultss


def find_web_to_text_div_and_p(url):
    dive_text = webpage_to_text_find_div_and_P(url)
    text = {'result': dive_text}
    return text


def webpage_to_text_p(url):
    url = url
    res = requests.get(url)
    html_page = res.content
    soup = BeautifulSoup(html_page, 'html.parser')
    result = soup.find_all("p")
    blacklist = [
        '[document]',
        'noscript',
        'header',
        'html',
        'meta',
        'head',
        'input',
        'script',
        # there may be more elements you don't want, such as "style", etc.
    ]
    resultss = []
    """Remove BlackListed Tag And extract Text"""
    for t in result:
        if t.parent.name not in blacklist:
            text = str(t.text)
            """Process For Cleaning The Text"""
            remove_n = text.replace('\n', ' ')
            clean_text = re.sub(' +', ' ', remove_n)
            remove_leading_space = clean_text.lstrip()
            remove_digits = str.maketrans('', '', digits)
            digit_removed = remove_leading_space.translate(remove_digits)

            if not {t.parent.name: remove_leading_space} in resultss:
                if remove_leading_space and t.parent.name == 'div':
                    resultss.append({t.parent.name: remove_leading_space})
            """SEt Tag name Of particular Text"""
            if not {'p': digit_removed} in resultss:
                if digit_removed:
                    resultss.append({'p': digit_removed})
    return resultss


def webpage_to_text_ul(url):
    url = url
    res = requests.get(url)
    html_page = res.content
    soup = BeautifulSoup(html_page, 'html.parser')
    result = soup.find_all("ul")
    resultss = []
    """Remove BlackListed Tag And extract Text"""
    for t in result:
        text = str(t.text)
        """Process For Cleaning The Text"""
        remove_n = text.replace('\n', ' ')
        clean_text = re.sub(' +', ' ', remove_n)
        remove_leading_space = clean_text.lstrip()
        remove_digits = str.maketrans('', '', digits)
        digit_removed = remove_leading_space.translate(remove_digits)
        alphanumeric = [character for character in digit_removed if
                        character.isalnum() or character.isspace() or character.isalpha()]
        alphanumeric = "".join(alphanumeric)

        """SEt Tag name Of particular Text"""
        if not {'ul': digit_removed} in resultss:
            if digit_removed:
                resultss.append({'ul': digit_removed})
    return resultss


def webpage_to_text_li(url):
    url = url
    res = requests.get(url)
    html_page = res.content
    soup = BeautifulSoup(html_page, 'html.parser')
    result = soup.find_all("li")
    resultss = []
    """Remove BlackListed Tag And extract Text"""
    for t in result:
        text = str(t.text)
        """Process For Cleaning The Text"""
        remove_n = text.replace('\n', ' ')
        clean_text = re.sub(' +', ' ', remove_n)
        remove_leading_space = clean_text.lstrip()
        remove_digits = str.maketrans('', '', digits)
        digit_removed = remove_leading_space.translate(remove_digits)
        """SEt Tag name Of particular Text"""
        if not {'li': digit_removed} in resultss:
            if digit_removed:
                resultss.append({'li': digit_removed})
    return resultss


"""End Of Flow"""


# """Use Docker support For Scraping"""
def Scrap_google_result(keywords, num_pages):
    # list = ["duckduckgo",  "bing", "google"]
    list = ["bing", "google"]
    crawled_data = []
    for i in range(len(list)):
        file_type = str(list[i])
        scrap_info = {
            "browser_config": {
                "random_user_agent": True
            },
            "scrape_config": {
                "search_engine": file_type,
                "keywords": keywords,
                "num_pages": num_pages
            }
        }
        resp = requests.post('http://localhost:3000', json=scrap_info)
        data = resp.json()
        crawled_data.append(data)
    links = []
    for data in crawled_data:
        for key_value, data in values.items():
            links.append(data['results'])
        for keywords, values in data['results'].items():
            pass
    seen = set()
    final_result = []
    for datas in links:
        for data in datas:
            if data['link'] not in seen:
                final_result.append(data)
                seen.add(data['link'])
            else:
                pass
    return final_result


# """Use Node_js Api For scraping"""
def multi_processing_Scrap_web_result(keywords, num_pages, search_engine):
    crawled_data = None
    file_type = str(search_engine)
    scrap_info = {
        'keyword': keywords,
        'num_pages': num_pages,
        'search_engine': file_type
    }
    resp = requests.post(settings.google_sc_scraper_api, json=scrap_info)
    data = resp.json()
    for key, urls in data['scraper_result'].items():
        if not urls['no_results']:
            if crawled_data is None:
                crawled_data = urls['results']
            else:
                crawled_data += urls['results']
    crawled_data = pd.DataFrame(crawled_data).drop_duplicates("link").to_dict('records')
    return crawled_data


def Scrap_web_result(keywords, num_pages):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        p1_google_result = executor.submit(multi_processing_Scrap_web_result, keywords, num_pages, 'google')
        p2_bing_result = executor.submit(multi_processing_Scrap_web_result, keywords, num_pages, 'bing')
        google_bing_result = p1_google_result.result() + p2_bing_result.result()
    seen = set()
    crawled_data = []
    for dict in google_bing_result:
        links = tuple(dict.items())
        if links[0] not in seen:
            seen.add(links[0])
            crawled_data.append(dict)
            print("returning crowled data...")
    return crawled_data

def Scrap_web_result_multisearch(keywords, num_pages):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        p0_csv_result = []
        with open('media/Dataset/domain_suffix.csv', mode='r') as csv_file:
            for line in csv_file.readlines():
                    data = line.partition(',')[-1].split()[0]
                    p0_csv_result = data
                    print("******",p0_csv_result)
                    if p0_csv_result:
                        return p0_csv_result 

        p1_google_result = executor.submit(multi_processing_Scrap_web_result, keywords, num_pages, 'google')
        p2_bing_result = executor.submit(multi_processing_Scrap_web_result, keywords, num_pages, 'google')
        google_bing_result = p1_google_result.result() + p2_bing_result.result() + p0_csv_result

    seen = set()
    crawled_data = []
    for dict in google_bing_result:
        links = tuple(dict.items())
        if links[0] not in seen:
            seen.add(links[0])
            crawled_data.append(dict)
    return crawled_data

def Google_scraper(keywords, num_pages):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        p1_google_result = executor.submit(multi_processing_Scrap_web_result, keywords, num_pages, 'google')
        google_result = p1_google_result.result() 
    seen = set()
    crawled_data = []
    for dict in google_result:
        links = tuple(dict.items())
        if links[0] not in seen:
            seen.add(links[0])
            crawled_data.append(dict)
    return crawled_data

def find_word_similarity(text1, text2):
    word_value = text_to_Words(text1, text2)
    return word_value

def Bing_scraper(keywords, num_pages):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        p1_bing_result = executor.submit(multi_processing_Scrap_web_result, keywords, num_pages, 'google')
       
        bing_result = p1_bing_result.result()
    seen = set()
    crawled_data = []
    for dict in bing_result:
        links = tuple(dict.items())
        if links[0] not in seen:
            seen.add(links[0])
            crawled_data.append(dict)
    return crawled_data

def Website_crawler(Keywords):
    keyword = Keywords[0]
    # Create and configure logger
    logging.basicConfig(filename="crawler_logs/crawling.log",
                        format='%(asctime)s %(message)s',
                        filemode='w')

    # Creating an object
    logger = logging.getLogger()

    # Setting the threshold of logger to DEBUG
    logger.setLevel(logging.DEBUG)

    keyword_searchId = KeywordSearch.objects.get(pk=keyword['id'])
    keyword_searchId.status = 1
    keyword_searchId.save()
    keyword_categories = Keyword_category.objects.filter(keywordId=keyword['id']).values_list('name',
                                                                                              flat=True).distinct()
    try:
        """Scraper function TO extract records from Google and bing"""
        crawled_data = Google_scraper([keyword['keyWord']], 25)
        keyword_searchId.total_crawled_result = len(crawled_data)
        keyword_searchId.save()
    except MemoryError as error:
        logger.info("MemoryError")
        logger.info(error)
    except Exception as exception:
        logger.info("somthing Wrong with Node_api: ")
        logger.info(exception)
        print('error')
    """Start Multi-Processing with links"""
    try:
        print(crawled_data)
        count_total_processed_links = 0
        logger.info(len(crawled_data))
        with concurrent.futures.ProcessPoolExecutor(max_workers=3) as executor:  # max_workers=1
            results = [
                executor.submit(multi_processing_crawl, store_link, keyword_searchId, keyword_categories, keyword) for
                store_link in crawled_data]
            for f in concurrent.futures.as_completed(results):
                print(f.result())
                logger.info(f.result())
                count_total_processed_links += 1
                logger.info(count_total_processed_links)
                print(len(crawled_data))
                print(count_total_processed_links)
                if len(crawled_data) - 1 == count_total_processed_links:
                    keyword_searchId.status = 2
                    keyword_searchId.save()
                    logger.info("Keyword Crawled Successfully")
                    return 'Success'
        keyword_searchId.status = 2
        keyword_searchId.save()
        logger.info("Keyword Crawled Successfully")
        """Restarting Sub-Process after compliting"""
        return 'Success'
    except MemoryError as error:
        logger.info("MemoryError On Multi-Processing")
        logger.info(error)

    except Exception as exception:
        print(exception)
        logger.info("exception On Multi-Processing")
        logger.info(exception)
        return 'Error'

def multi_processing_crawl(store_link, keyword_searchId, keyword_categories, keyword):
    try:
        # print('url_coming_for_multi_processing')
        
        """Convert Website_Link to Text"""
        extracted_text = url_to_text(store_link['link'])
        if extracted_text == '':
            pass
        """Check If description is not null"""
        if keyword['description'] != '':
            """Find Similarity / Rating """
            similarity_dict = find_cosine_similarity_with_keyword_repetition(keyword['description'],
                                                                             str(extracted_text)) 

            if similarity_dict['total_percentage'] > 20:
                similarity = similarity_dict['total_percentage']

        elif keyword['description'] == '':
            similarity = 0
        """Classify the text an get category of text"""
        # type_of_text = text_classification_by_repeated_keywords({'text': str(extracted_text['result'])})
        type_of_text = find_classification([store_link['link']])

        """Separate sub_category by it's domain"""
        info = tldextract.extract(store_link['link'])
        if info.suffix == 'com':
            url_extension_type = 'COMMERCIAL'
        elif info.suffix == 'org':
            url_extension_type = 'ORGANIZATION'
        elif info.suffix == 'edu':
            url_extension_type = 'EDUCATIONAL'
        else:
            url_extension_type = 'OTHERS'

        """Overriding url_extension_type value if sub_category in not null"""
        print(type_of_text[0])

        if 'sub_category' in type_of_text[0]:
            if type_of_text[0]['sub_category'] != '':
                url_extension_type = type_of_text[0]['sub_category']
        """Check sub_category is selected or not if yes than stored in db"""
        print(type_of_text[0]['category'])
        if type_of_text[0]['category'] in keyword_categories:
            """Check if Url already stored in DB with same keyword_id"""
            SearchResults = SearchResult.objects.filter(
                Q(title=store_link['title']) & Q(keywordId=keyword_searchId) |
                Q(keywordId=keyword_searchId) & Q(url=store_link['link'])).values()
        
            if len(SearchResults) == 0:
                result_data = SearchResult.objects.create(title=store_link['title'],
                                                          url=store_link['link'],
                                                          description=store_link['snippet'],
                                                          keywordId=keyword_searchId,
                                                          url_extension='.' + info.suffix,
                                                          type_of_text=type_of_text[0]['category'],
                                                          url_extension_type=url_extension_type,
                                                          matched_similarity=similarity)
                """Store Company_Info Of Particular Url"""
                if type_of_text[0]['category'] == "BUSINESS":
                    try:
                        company_details = get_companies_details_from_googlemaps(
                            {'url': store_link['link']})
                        print(company_details)
                    except Exception as exception:
                        print(exception)
                        pass
                    try:
                        if len(company_details) >= 1:
                            store_result_id = SearchResult.objects.get(id=result_data.id)
                            is_exist = company_info.objects.filter(keywordId=keyword_searchId,
                                                                   company_name=company_details[0][
                                                                       'name']).values()
                            if len(is_exist) == 0:
                                company_info.objects.create(company_name=company_details[0]['name'],
                                                            icon=company_details[0]['icon'],
                                                            place_id=company_details[0]['place_id'],
                                                            rating=company_details[0]['rating'],
                                                            keywordId=keyword_searchId,
                                                            searchResultId=store_result_id,
                                                            company_url=store_link['link'],
                                                            matched_similarity=similarity)
                    except Exception as exception:
                        print(exception)
                        return '{ url : ' + store_link['link'] + '} Exception occurs For Google map api'
            """Return Massage For log FIle"""
            return '{ url : ' + store_link['link'] + '} Stored in Database successfully '
        return '{ url : ' + store_link['link'] + '} Category doesnt match to selected categories'
    except MemoryError as error:
        return "MemoryError On Process-and-saving-data"

    except Exception as exception:
        return 'Exception_Caught For : ' + 'url :' + store_link['link']

"""Start Text Comparision Utility Functions"""


def Text_Comparision_controller(all_text):
    text1 = all_text['text1'].lower()
    text2 = all_text['text2'].lower()
    search_keyword = all_text['search_keyword'].lower()
    keywords = keyword_counter(search_keyword, text2)
    vector1 = text_to_cosine(text1)
    vector2 = text_to_cosine(text2)
    similar_word_count = text_to_Words(text1, text2)
    cosine = get_cosine(vector1, vector2)
    value = cosine * 100

    m = difflib.SequenceMatcher(None, text1, text2)
    ratio = m.ratio() * 100
    final_result = (value + ratio + similar_word_count + keywords / 4)  # + keywords or + less_weightage
    if final_result >= 100:
        final_result = 100
    all_text['comparison_ratio'] = final_result
    return all_text


def get_cosine(vec1, vec2):
    intersection = set(vec1.keys()) & set(vec2.keys())
    numerator = sum([vec1[x] * vec2[x] for x in intersection])
    sum1 = sum([vec1[x] ** 2 for x in vec1.keys()])
    sum2 = sum([vec2[x] ** 2 for x in vec2.keys()])
    denominator = math.sqrt(sum1) * math.sqrt(sum2)

    if not denominator:
        return 0.0
    else:
        return float(numerator) / denominator


def text_to_cosine(text):
    search = str(text)
    WORD = re.compile(r'\w+')
    words = WORD.findall(search)
    return Counter(words)


def text_to_Words(text, text2):
    text1_tokens = word_tokenize(text)
    text2_tokens = word_tokenize(text2)
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]
    setA = set(tokens_without_sw_text1)
    setB = set(tokens_without_sw_text2)

    overlap = setA & setB
    universe = setA | setB
    result3 = float(len(overlap)) / len(universe) * 100
    return result3


"""End Of replica"""


def keyword_counter(keyword, text2):
    count = text2.count(keyword)
    return count


"""End Of Text Comparision Utility Functions"""


def url_classification(url):
    if url.startswith('https://'):
        url = url.replace("https://", "")
    elif url.startswith('http://'):
        url = url.replace("http://", "")

    """Runniong Three Functions parallel"""
    info = tldextract.extract(url)
    test = info.domain
    category = []
    with open('media/Dataset/domain_suffix.csv', mode='r') as csv_file:
        for line in csv_file.readlines():
            if test in line:
                data = line.partition(',')[-1].split()[0]
                category = data
    if category:
        return category

    names = ['URL', 'Category']
    df = pd.read_csv('media/Dataset/url_classification.csv', names=names, na_filter=False)

    X_train = df['URL']
    y_train = df['Category']
    X_train.shape
    text_clf = Pipeline([('vect', CountVectorizer()), ('tfidf', TfidfTransformer()), ('clf', MultinomialNB())])
    text_clf = text_clf.fit(X_train, y_train)
    n_iter_search = 5
    parameters = {'vect__ngram_range': [(1, 1), (1, 2)], 'tfidf__use_idf': (True, False), 'clf__alpha': (1e-2, 1e-3)}
    gs_clf = RandomizedSearchCV(text_clf, parameters, n_iter=n_iter_search)
    gs_clf = gs_clf.fit(X_train, y_train)
    url = gs_clf.predict([url])
    if url == ['TECH']:
        url = ['TECHNOLOGY']
    return url[0]


def find_weburl_classification_no_autoclassification(url):
    if url.startswith('https://'):
        url = url.replace("https://", "")
    elif url.startswith('http://'):
        url = url.replace("http://", "")
    names = ['URL', 'Category']
    df = pd.read_csv('media/Dataset/url_classification.csv', names=names, na_filter=False)

    X_train = df['URL']
    y_train = df['Category']
    X_train.shape
    text_clf = Pipeline([('vect', CountVectorizer()), ('tfidf', TfidfTransformer()), ('clf', MultinomialNB())])
    text_clf = text_clf.fit(X_train, y_train)
    n_iter_search = 5
    parameters = {'vect__ngram_range': [(1, 1), (1, 2)], 'tfidf__use_idf': (True, False), 'clf__alpha': (1e-2, 1e-3)}
    gs_clf = RandomizedSearchCV(text_clf, parameters, n_iter=n_iter_search)
    gs_clf = gs_clf.fit(X_train, y_train)
    url = gs_clf.predict([url])
    if url == ['TECH']:
        url = ['TECHNOLOGY']
    return url[0]


def find_web_text_classification(text_or_url):
    if 'text' in text_or_url:
        text = text_or_url['text']
    elif 'url' in text_or_url:
        url = text_or_url['url']
        text = webpage_to_text_find_div_and_P(url)
    train = [
        (
            'Companies, Company, Startup, Startups, Manufacturers, Wholesalers, Affiliate Merchants, Franchisors, Importers,'
            'Dropshippers, Manufacturer, Vendors, Vendor, Wholesaler, Distributor, Affiliate Merchant, Franchisor, Importer,'
            'Exporters, Exporter, Social Advocacy Groups, Charitable Organizations, Foundations, Civic League, Trade Association,'
            'Social Club, Fraternal Societies, Employee Beneficiary Association, Domestic Fraternal Societies, Social Welfare Organization,'
            'Local Employee Association, Domestic Fraternal Associations, Associations, Association, Recreational Club, Professional Association,'
            'Organization, Organizations, Government, Venture Capital Firms, Family Office, Angel groups, Banks Agencies, Angel Investor,'
            'Accelerators, Corporate Investors, Incubators, Government Agencies, Government Agency, Banks Agency, Agency, Agencies,'
            'start-up, start-ups, enterprise, commercial, corporate, investment, profitability, industry, industries, founder,'
            'founders, entrepreneur, ventures, corporation, board of director, team, portfolio, services, demo, client, customer,'
            'firms, sector, strategic, employment, office, job, jobs, about us, contact us, brand, industrial, professional, sale.'
            'entity, contractors, partnership, partnerships, contractor, dealers, dealer, capital, headquarter, headquarters,'
            'facility',
            'BUSINESS'),
        (
            'Patent, Scientific Papers, Innovation, Technology, Science, Engineer, Software, Hardware, Blockchain, Artificial Intelligence,'
            'Virtual Reality, Augmented Reality, AI, VR, Technology description, tech, science and technology, new technology,'
            'science technology, Technical, technological, mechanical, engineering, scientific, engineer, physical, analytical, fundamental,'
            'feasability, study, computerized, hi-tech, nuclear, medical, green energy, agricultural, electronic, innovative, optical,'
            'laser, mobile, smart, multimedia, biotechnology, computing, informatics, computer science,'
            'researchcyber, cybersecutiry, computerization',
            'TECHNOLOGY'),
        (
            'market, trend, financial, forecast, Market Overview, CAGR, Segmentation, Regional Analysis, Analysis, Analyses,Global market,'
            'Region, MRFR, Competitors, Key Competitors, Report, Market Segmentation, Grow, Growth, Growth rate, rate, market value,'
            'year, MARKET LANDSCAPE, Market ecosystem, landscape, Market characteristics, Market segmentation analysis, Market definition,'
            'MARKET SIZING, Market size, Markek forecast, FIVE FORCES ANALYSIS. Bargaining power of buyers, Bargaining power of suppliers,'
            'Threat of new entrants, Threat of substitutes, Threat of rivalry, Market condition, MARKET SEGMENTATION BY PRODUCT, Segmentation by product,'
            'Comparison by product, Market size and forecast, Market opportunity, Market opportunity by product, CUSTOMER LANDSCAPE, REGIONAL LANDSCAPE,'
            'Geographical segmentation, Regional comparison, Americas, EMEA, APAC, Key leading countries, Market opportunity, DECISION FRAMEWORK,'
            'DRIVERS AND CHALLENGES, Market drivers, Market challenges, MARKET TRENDS, Worldwide Opportunities, Increasing adoption, Driving Trends,'
            'Growing demand, Future Potential Growth, stock, VENDOR LANDSCAPE, economic, Landscape disruption, VENDOR ANALYSIS, Vendors covered, Vendor classification,'
            'Market positioning of vendors, North America, Asia Pacific, Middle East, Africa, Europe, South America',
            'FINANCIAL'),
        (
            'Articles, Video, Audio, Journalists, Print Media, Newspaper, National Newspapers, Daily newspapers, Special Audience newspapers,'
            'Magazines, Consumer magazines, General Interest Magazine, Glamour Magazine, Film Magazine, Special Interest Magazine, Business Publications,'
            'Business Publications, Professional Publications, Trade Journals, Industrial Publications, Broadcast Media, Television, Radio,'
            'Internet, social networking, blog, news, informational, Institutional Publications, Lifestyle Magazine, post, tech news, latest technology, technews,'
            'latest technology news, cool tech, technology review, newtech, technology articles, news technology, technology updates, business news,'
            'mass media, press, journalists, newspapers, television, broadcasters, broadcasting, journalism, reporters, broadcast,'
            'multimedia, communications, journalistic, newsworthy, pressroom, newsroom, mainstream, tabloid, entertainment, propaganda',
            'MEDIA')
    ]
    cl = NaiveBayesClassifier(train)
    classify_text = cl.classify(text)
    return classify_text


"""Get Company Using GoogleMap"""


def get_companies_details_from_googlemaps(info):
    api_key = ''

    # url variable store url
    url = "https://maps.googleapis.com/maps/api/place/textsearch/json?"
    # The text string on which to search

    if 'location' in info:
        company_name = re.sub('\s+', '+', info['name'])
        company_location = re.sub('\s+', '+', info['location'])
        query = company_name + ',' + company_location
    elif 'url' in info:
        domain = re.sub('\s+', '+', info['url'])
        info = tldextract.extract(domain)
        domain_name = info.domain
        query = domain_name
    elif 'name' in info:
        company_name = re.sub('\s+', '+', info['name'])
        query = company_name
    """"get method of requests module
        return response object"""
    r = requests.get(url + 'query=' + query + '&radius=50000&key=' + api_key)
    #  json format data into python format data
    x = r.json()

    """now x contains list of nested dictionaries
       we know dictionary contain key value pair
       store the value of result key in variable y"""
    y = x['results']
    return y


def get_from_google_maps(info):
    api_key = ''

    # url variable store url
    url = "https://maps.googleapis.com/maps/api/place/textsearch/json?"
    # The text string on which to search
    company_name = re.sub('\s+', '+', info['name'])
    company_location = re.sub('\s+', '+', info['location'])
    query = company_name + ',' + company_location
    """"get method of requests module
        return response object"""
    r = requests.get(url + 'query=' + query + '&radius=50000&key=' + api_key)
    #  json format data into python format data
    x = r.json()

    """now x contains list of nested dictionaries
       we know dictionary contain key value pair
       store the value of result key in variable y"""
    y = x['results']
    return y


"""FIND COSINE Similarity"""


def find_cosine_similarity(text1, text2):
    vector1 = text_to_cosine(text1)
    vector2 = text_to_cosine(text2)
    cosine = get_cosine(vector1, vector2)
    Cosine_value = cosine * 100
    return Cosine_value

""" replica """

def replica_cosine_similarity(text1, text2):
    val = []
    new_lines = []
    phrases = []
    for line in text2:
        line = line.strip()
        if line not in new_lines:
            new_lines.append(line)
    print(new_lines)
    lines = list(set(text2))
    for line in new_lines:
        vector1 = text_to_cosine(text1)
        vector2 = text_to_cosine(line)
        cosine = get_cosine(vector1, vector2)
        Cosine_value = cosine * 100
        if Cosine_value != 0.0:
            value = {Cosine_value}
            linecoisne = {line:value}
            phrases.append(linecoisne)
            val.append(value)
    number = len(val)
    percentage =[]
    for vals in val:
        word_val = sum(vals)
        percentage.append(word_val)
    result = sum(percentage)/number
    data = {"percentage":result, "phrases" : phrases }
    return data

def find_uniquePhrases(words, text2):
    sentences = text2.split('.')
    lines = []
    for sentence in sentences :
         for word in words :
               if word in sentence :
                    line = sentence
                    lines.append(line)
    
    return lines

"""Similarity_keyword return % value"""


def find_similarity_keyword(text, text2):
    """For total Words of Test2"""
    split_list_text1 = text.split()
    split_list_text2 = text2.split()

    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text.lower())
    text2_tokens = word_tokenize(text2.lower())

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [wnl.lemmatize(i) for i in tokens_without_sw_text1]
    tokens_without_sw_text2 = [wnl.lemmatize(i) for i in tokens_without_sw_text2]

    """Removing unnecessary single line keywords"""
    tokens_without_sw_text1 = [i for i in tokens_without_sw_text1 if len(i) > 1]
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]
    """Removing declension(ed) from the text1 and overriding the value"""
    for text1_words in tokens_without_sw_text1:
        for suffix in ['ed']:
            if text1_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text1.index(text1_words)
                """Pop that word from list"""
                tokens_without_sw_text1.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text1_words[:-len(
                    suffix)]
                tokens_without_sw_text1.append(result)

    """Removing declension(ed) from the text1 and overriding the value"""
    for text2_words in tokens_without_sw_text2:
        for suffix in ['ed']:
            if text2_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text2.index(text2_words)
                """Pop that word from list"""
                tokens_without_sw_text2.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text2_words[:-len(suffix)]
                tokens_without_sw_text2.append(result)

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Creating "set" of both the keyword list"""
    setA = set(alphanumeric_for_text1)
    setB = set(alphanumeric_for_text2)

    """Separating the similar and non similar words From both the list of text"""
    overlap = setA & setB
    universe = setA | setB
    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)
    """Creating set For list of text1(FOR_: tot_unique_k)"""
    set_counter_text1 = set(counts_for_text1)

    """Counting repeated words(FOR_: tot_k_rep)"""
    count_rep_keyword_text1 = 0
    for key, count in counts_for_text1.items():
        count_rep_keyword_text1 += count

    """Calculating the repeated words in list of text2"""
    counts_for_text2 = Counter(alphanumeric_for_text2)
    set_counter = set(counts_for_text2)
    """Separating similar key_words from both the list of text"""
    similar_words = set_counter_text1 & set_counter
    tot_key_rep_text2 = 0
    tot_rep_key = 0
    """Count total repeated key_word from text1 to text2(FOR_: tot_ref_k) and total key_words(FOR_: repetition)"""
    for key, count in counts_for_text2.items():
        tot_key_rep_text2 += count
        if key in counts_for_text1:
            tot_rep_key += count  # + counts_for_text1[key]
    """Adding Total number of keywords and list of keywords"""
    unique_key_text1 = len(counts_for_text1)
    unique_key_text2 = len(counts_for_text2)
    list_keyword_text1 = counts_for_text1
    list_keyword_text2 = counts_for_text2

    """Adding Total number of keywords_repeated"""
    tot_key_rep_text1 = count_rep_keyword_text1
    tot_key_rep_text2 = tot_key_rep_text2

    """Adding Total number of keywords are present in text2 and list of those keywords"""
    tot_ref_key = len(overlap)
    list_ref_key = {}
    total_list_ref_count = 0
    for value in overlap:
        total_count = counts_for_text1[value] + counts_for_text2[value]
        list_ref_key[value] = total_count
        total_list_ref_count += total_count
    list_ref_key['total_count'] = total_list_ref_count

    """Set dictionary To descending order"""
    list_ref_key = dict(sorted(list_ref_key.items(), key=operator.itemgetter(1), reverse=True))
    list_keyword_text1 = dict(sorted(list_keyword_text1.items(), key=operator.itemgetter(1), reverse=True))
    list_keyword_text2 = dict(sorted(list_keyword_text2.items(), key=operator.itemgetter(1), reverse=True))

    """Calculating Percentage For both the list of Text"""
    result = float(len(overlap)) / len(universe) * 100
    data = {'percentage': result,
            'unique_key_text1': unique_key_text1,
            'unique_key_text2': unique_key_text2,
            'list_keyword_text1': list_keyword_text1,
            'list_keyword_text2': list_keyword_text2,
            'tot_key_rep_text1': tot_key_rep_text1,
            'tot_key_rep_text2': tot_key_rep_text2,
            'tot_ref_key': tot_ref_key,
            'list_ref_key': list_ref_key,
            'tot_rep_key': tot_rep_key,
            'total_word_text1': len(split_list_text1),
            'total_word_text2': len(split_list_text2)
            }
    return data

""" replica for purpose task """

def replica_similarity_keyword(text, text2):
    """For total Words of Test2"""
    split_list_text1 = text.split()
    split_list_text2 = text2.split()

    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text.lower())
    text2_tokens = word_tokenize(text2.lower())

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [wnl.lemmatize(i) for i in tokens_without_sw_text1]
    tokens_without_sw_text2 = [wnl.lemmatize(i) for i in tokens_without_sw_text2]

    """Removing unnecessary single line keywords"""
    tokens_without_sw_text1 = [i for i in tokens_without_sw_text1 if len(i) > 1]
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]
    """Removing declension(ed) from the text1 and overriding the value"""
    for text1_words in tokens_without_sw_text1:
        for suffix in ['ed']:
            if text1_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text1.index(text1_words)
                """Pop that word from list"""
                tokens_without_sw_text1.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text1_words[:-len(
                    suffix)]
                tokens_without_sw_text1.append(result)

    """Removing declension(ed) from the text1 and overriding the value"""
    for text2_words in tokens_without_sw_text2:
        for suffix in ['ed']:
            if text2_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text2.index(text2_words)
                """Pop that word from list"""
                tokens_without_sw_text2.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text2_words[:-len(suffix)]
                tokens_without_sw_text2.append(result)

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Creating "set" of both the keyword list"""
    setA = set(alphanumeric_for_text1)
    setB = set(alphanumeric_for_text2)

    """Separating the similar and non similar words From both the list of text"""
    overlap = setA & setB
    universe = setA | setB
    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)
    """Creating set For list of text1(FOR_: tot_unique_k)"""
    set_counter_text1 = set(counts_for_text1)
    tot_ref_key = len(overlap)
    return set_counter_text1  


def find_keyword_repetition(text, text2):
    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text.lower())
    text2_tokens = word_tokenize(text2.lower())

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [wnl.lemmatize(i) for i in tokens_without_sw_text1]
    tokens_without_sw_text2 = [wnl.lemmatize(i) for i in tokens_without_sw_text2]
    """Removing unnecessary single line keywords"""
    tokens_without_sw_text1 = [i for i in tokens_without_sw_text1 if len(i) > 1]
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]

    """Removing declension(ed) from the text1 and overriding the value"""
    for text1_words in tokens_without_sw_text1:
        for suffix in ['ed']:
            if text1_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text1.index(text1_words)
                """Pop that word from list"""
                tokens_without_sw_text1.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text1_words[:-len(suffix)]
                tokens_without_sw_text1.append(result)

    """Removing declension(ed) from the text1 and overriding the value"""
    for text2_words in tokens_without_sw_text2:
        for suffix in ['ed']:
            if text2_words.endswith(suffix):
                """Taking index value of "ed" word"""
                index = tokens_without_sw_text2.index(text2_words)
                """Pop that word from list"""
                tokens_without_sw_text2.pop(index)
                """Removing "ed" from the word and append into list"""
                result = text2_words[:-len(suffix)]
                tokens_without_sw_text2.append(result)

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)

    """Calculating the repeated words in list of text2"""
    counts_for_text2 = Counter(alphanumeric_for_text2)

    tot_key_rep_text2 = 0
    tot_rep_key = 0
    """Count total repeated key_word from text1 to text2(FOR_: tot_ref_k) and total key_words(FOR_: repetition)"""
    for key, count in counts_for_text2.items():
        tot_key_rep_text2 += count
        if key in counts_for_text1:
            tot_rep_key += count  # + counts_for_text1[key]

    """Adding Total number of keywords_repeated"""
    tot_key_rep_text2 = tot_key_rep_text2

    """Calculating Percentage For both the list of Text"""
    result = tot_rep_key / tot_key_rep_text2 * 100

    data = {'percentage': result}
    return data


def find_cosine_similarity_with_keyword_repetition(text, text2):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        """Find cosine_similarity and keyword repetition value parallel"""
        cosine_value = executor.submit(find_cosine_similarity, text, text2)
        keyword_repetition_value = executor.submit(find_keyword_repetition, text, text2)

        """Take Result From cosine and keyword repetition"""
        cosine_value = cosine_value.result()
        keyword_repetition_value = keyword_repetition_value.result()

        """Find Rating value"""
        total_percentage = find_rating([cosine_value], [keyword_repetition_value['percentage']])
        return {'total_percentage': total_percentage['Total_Rating'][0]['rating_0'],
                'cosine_value': int(cosine_value),
                'keyword_repetition_value': int(keyword_repetition_value['percentage'])}


def find_rating(cosine_value, keyword_repetition_value):
    if len(cosine_value) != len(keyword_repetition_value):
        return {"Error": 'Value is missing'}
    rating_result = []
    for index_value in range(len(cosine_value)):
        calculating_formula = int(cosine_value[index_value]) * 75 + int(keyword_repetition_value[index_value]) * 25
        calculating_parentage = int(calculating_formula / 100)
        if int(keyword_repetition_value[index_value]) >= 100:
            calculating_parentage = int(calculating_parentage) + int(calculating_parentage) * 25 / 100
        if int(calculating_parentage) >= 100:
            calculating_parentage = 100
        rating_result.append({'rating_' + str(index_value): int(calculating_parentage)})
    return {"Total_Rating": rating_result}


def webpage_to_text_find_div_and_P(url):
    url = url
    res = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
    html_page = res.content
    soup = BeautifulSoup(html_page, 'html.parser')
    result = soup.find_all("div")
    blacklist = [
        '[document]',
        'noscript',
        'header',
        'html',
        'meta',
        'head',
        'input',
        'script',
    ]
    data = ''

    for t in result:
        if t.find('p'):
            text = t.find('p').text
            remove_n = text.replace('\n', ' ')
            remove_r = remove_n.replace('\r', ' ')
            remove_t = remove_r.replace('\t', ' ')
            clean_text = re.sub(' +', ' ', remove_t)
            remove_leading_space = clean_text.lstrip()
            remove_digits = str.maketrans('', '', digits)
            digit_removed = remove_leading_space.translate(remove_digits)
            if not digit_removed in data:
                if digit_removed:
                    data += digit_removed

        if t.parent.name not in blacklist:
            text = str(t.text)
            remove_n = text.replace('\n', ' ')
            remove_r = remove_n.replace('\r', ' ')
            remove_t = remove_r.replace('\t', ' ')
            clean_text = re.sub(' +', ' ', remove_t)
            remove_leading_space = remove_t.lstrip()
            if not remove_leading_space in data:
                if remove_leading_space and t.parent.name == 'div':
                    data += remove_leading_space

    return data


"""New classification Function For categorise by max repeated Keywords"""


def repeated_value_media(MEDIA, text):
    """Calling Function which can Find TOTAL-Repeated value"""
    parentage = find_total_repetition(MEDIA, text)
    return parentage['tot_rep_key']


def repeated_value_technology(TECHNOLOGY, text):
    """Calling Function which can Find TOTAL-Repeated value"""
    parentage = find_total_repetition(TECHNOLOGY, text)
    return parentage['tot_rep_key']


def repeated_value_financial(FINANCIAL, text):
    """Calling Function which can Find TOTAL-Repeated value"""
    parentage = find_total_repetition(FINANCIAL, text)
    return parentage['tot_rep_key']


def repeated_value_business(BUSINESS, text):
    """Calling Function which can Find TOTAL-Repeated value"""
    parentage = find_total_repetition(BUSINESS, text)
    return parentage['tot_rep_key']


def text_classification_by_repeated_keywords(text_or_url):
    if 'text' in text_or_url:
        if text_or_url['text'] == '':
            return "OTHERS"
        """Take text For find Category"""
        text2 = text_or_url['text']
    else:
        extracted_text = url_to_text(text_or_url['url'])
        if extracted_text == '':
            return "OTHERS"
        text2 = extracted_text
    TECHNOLOGY = ''
    MEDIA = ''
    FINANCIAL = ''
    BUSINESS = ''
    """Read Text From csv files and separate keywords"""
    with open('media/Dataset/text_classification.csv', mode='r') as csv_file:
        csv_reader = csv.DictReader(csv_file)
        line_count = 0
        for row in csv_reader:
            if line_count == 0:
                line_count += 1
            if row["TECHNOLOGY"]:
                TECHNOLOGY += row["TECHNOLOGY"] + '\n'
            if row["MEDIA"]:
                MEDIA += row["MEDIA"] + '\n'
            if row["FINANCIAL"]:
                FINANCIAL += row["FINANCIAL"] + '\n'
            if row["BUSINESS"]:
                BUSINESS += row["BUSINESS"] + '\n'

            line_count += 1

    """Tokenize The text and extract keywords"""
    text2_tokens = word_tokenize(text2.lower())
    tokens_without_sw_text2 = [word for word in text2_tokens if not word in stopwords.words()]
    tokens_without_sw_text2 = [wnl.lemmatize(i) for i in tokens_without_sw_text2]
    """Removing unnecessary single line keywords"""
    tokens_without_sw_text2 = [i for i in tokens_without_sw_text2 if len(i) > 1]
    alphanumeric_for_text2 = [character for character in tokens_without_sw_text2 if
                              character.isalnum() or character.isspace() or character.isalpha()]

    """Count Repeated keywords"""
    counts_for_text2 = Counter(alphanumeric_for_text2)
    list_s = {}

    """Start Multiprocessing to find repeated value to compare with  every dataset"""
    with concurrent.futures.ProcessPoolExecutor() as executor:
        """Calling all function parallel which find repeated value from - media, technology, financial, 
        business Dataset """
        p1 = executor.submit(repeated_value_media, MEDIA, counts_for_text2)
        p2 = executor.submit(repeated_value_technology, TECHNOLOGY, counts_for_text2)
        p3 = executor.submit(repeated_value_financial, FINANCIAL, counts_for_text2)
        p4 = executor.submit(repeated_value_business, BUSINESS, counts_for_text2)
        list_s['MEDIA'] = p1.result()
        list_s['TECHNOLOGY'] = p2.result()
        list_s['FINANCIAL'] = p3.result()
        list_s['BUSINESS'] = p4.result()
        if p1.result() == 0 and p2.result() == 0 and p3.result() == 0 and p4.result() == 0:
            return "OTHERS"

        """Take Highest Repeated value From all the Dataset"""
        category_by_MaxValue = max(list_s, key=list_s.get)
        return category_by_MaxValue


# FInd total Repeated Value from Two text
def find_total_repetition(text, counts_for_text2):
    """Tokenize Both the text"""
    text1_tokens = word_tokenize(text.lower())

    """Creating The Word List and removing stopwords(a, an ,the) Of all Tokenize text"""
    tokens_without_sw_text1 = [word for word in text1_tokens if not word in stopwords.words()]

    """Removing plural from both the text and overriding the value"""
    tokens_without_sw_text1 = [wnl.lemmatize(i) for i in tokens_without_sw_text1]

    """Removing wanted special_keywords(@#$%>,)"""
    alphanumeric_for_text1 = [character for character in tokens_without_sw_text1 if
                              character.isalnum() or character.isspace() or character.isalpha()]
    """Calculating the repeated words in list of text1"""
    counts_for_text1 = Counter(alphanumeric_for_text1)

    tot_rep_text2 = 0
    """Count total repeated key_word from text1 to text2(FOR_: tot_ref_k) and total key_words(FOR_: repetition)"""
    for key, count in counts_for_text2.items():
        if key in counts_for_text1:
            tot_rep_text2 += count  # + counts_for_text1[key]

    tot_rep_text1 = 0
    for key, count in counts_for_text1.items():
        if key in counts_for_text2:
            tot_rep_text1 += count  # + counts_for_text1[key]
    tot_rep_key = tot_rep_text2 + tot_rep_text1
    data = {'tot_rep_key': tot_rep_key}
    return data


"""End Of Flow"""

def find_classification(urls):
    classification = []
    with concurrent.futures.ProcessPoolExecutor(max_workers=1) as executor:
        """Call find_classification_multiProcess Function at sme time for all the links"""
        results = [
            executor.submit(find_classification_multiProcess, url) for
            url in urls]
        for category in concurrent.futures.as_completed(results):
            classification.append(category.result())
    return classification

def find_classification_multiProcess(url):
    parsed_uri = urlparse(url)
    domain = '{uri.netloc}/'.format(uri=parsed_uri)
    result = domain.replace('www.', '')

    category = []
    with open('media/Dataset/domain_suffix.csv', mode='r') as csv_file:
        for line in csv_file.readlines():
            if result in  line:
                data = line.partition(',')[-1].split()[0]
                try:
                    sub_category = data.partition(',')[-1].split()[0]
                    extract_category = data.replace(sub_category, '')
                    category = extract_category.replace(',', '')
                except:
                    category = data
                    sub_category = ''

    if len(category) != 0:  
        return {'url': url, 'category': category, 'sub_category': sub_category, 'web_url_result': 'prefixdatset',
                'web_text_clasification': 'prefixdatset'}
    with concurrent.futures.ProcessPoolExecutor() as executor:
        extracted_text_priority = executor.submit(url_to_text, url)
        extracted_text = extracted_text_priority.result()

        p2_web_text_clasification = executor.submit(text_classification_by_repeated_keywords,
                                                    {'text': extracted_text})
        web_text_clasification = p2_web_text_clasification.result()
        
        web_url_result = ''
        if web_text_clasification == 'BUSINESS' or web_text_clasification == 'MEDIA':
            web_url_result += web_text_clasification  
        else:
            p1_web_url_result = executor.submit(url_classification, url)
            data = p1_web_url_result.result()
            web_url_result =data.replace(',','')

        if web_text_clasification == web_url_result :
            category = web_text_clasification
        elif web_url_result == 'BUSINESS' or web_text_clasification == 'BUSINESS' and web_url_result != 'MEDIA' and web_text_clasification != 'MEDIA':
            category = 'BUSINESS'
        elif web_url_result == 'MEDIA' or web_text_clasification == 'MEDIA':
            category = 'MEDIA'
        elif web_url_result == 'MEDIA' and web_text_clasification == 'OTHERS':
            category = 'OTHERS'
        else:
            category = 'OTHERS'
    return {'url': url, 'category': category,"web_text_clasification": web_text_clasification,
            'web_url_result': web_url_result}

"""Replica For testing"""


def find_classification_no_auto(urls):
    classification = []
    with concurrent.futures.ProcessPoolExecutor(max_workers=1) as executor:
        """Call find_classification_multiProcess Function at sme time for all the links"""
        results = [
            executor.submit(find_classification_no_auto_multiProcess, url) for
            url in urls]
        for category in concurrent.futures.as_completed(results):
            classification.append(category.result())
    return classification


def find_classification_no_auto_multiProcess(url):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        extracted_text_priority = executor.submit(url_to_text, url)
        extracted_text = extracted_text_priority.result()

        p1_web_url_result = executor.submit(url_classification, url)
        p2_web_text_clasification = executor.submit(text_classification_by_repeated_keywords,
                                                    {'text': extracted_text})

        web_url_result = p1_web_url_result.result()
        web_text_clasification = p2_web_text_clasification.result()
        if web_url_result == web_text_clasification:
            category = web_url_result
        elif web_url_result == 'BUSINESS' or web_text_clasification == 'BUSINESS' and web_url_result != 'MEDIA' and web_text_clasification != 'MEDIA':
            category = 'BUSINESS'
        elif web_url_result == 'MEDIA' or web_text_clasification == 'MEDIA' and web_text_clasification != 'OTHERS':
            category = 'MEDIA'
        elif web_url_result == 'MEDIA' and web_text_clasification == 'OTHERS':
            category = 'OTHERS'
        else:
            category = 'OTHERS'
    return {'url': url, 'category': category, 'web_url_result': web_url_result,
            "web_text_clasification": web_text_clasification}


"""---END----"""
""" -- similarity -- """
def find_sentence_similarity(text1, text2):
    m = difflib.SequenceMatcher(None, text1, text2)
    ratio = m.ratio() * 100
    return ratio

# """find_rating_percentage Multi processing"""
def find_rating_percentage_list(urls, referance_text):
    all_ratings = []
    for url in urls:
        rating_result = find_rating_percentage_list_muliprocessing(url, referance_text)
        all_ratings.append(rating_result)
    return all_ratings

def find_rating_percentage_list_muliprocessing(url, referance_text):
    with concurrent.futures.ProcessPoolExecutor() as executor:
        p1_extracted_text = executor.submit(find_html_to_text, url)
        extracted_text = p1_extracted_text.result()
    similarity_dict = find_cosine_similarity_with_keyword_repetition(referance_text, str(extracted_text))
    similarity = similarity_dict['total_percentage']
    return {'link': url, 'similarity': similarity}

def find_html_to_text(url):
    response = requests.get(url, "lxml", headers={'User-Agent': 'Mozilla/5.0'})
    text = get_text(response.text)
    remove_n = text.replace('\n', ' ')
    remove_r = remove_n.replace('\r', ' ')
    remove_t = remove_r.replace('\t', ' ')
    clean_text = re.sub(' +', ' ', remove_t)
    remove_leading_space = clean_text.lstrip()
    return remove_leading_space


def text_url_classification(text):
    names = ['text', 'Category']
    df = pd.read_csv('media/Dataset/texturldataset.csv', names=names, na_filter=False)

    X_train = df['text']
    y_train = df['Category']
    X_train.shape
    text_clf = Pipeline([('vect', CountVectorizer()), ('tfidf', TfidfTransformer()), ('clf', MultinomialNB())])
    text_clf = text_clf.fit(X_train, y_train)
    n_iter_search = 5
    parameters = {'vect__ngram_range': [(1, 1), (1, 2)], 'tfidf__use_idf': (True, False), 'clf__alpha': (1e-2, 1e-3)}
    gs_clf = RandomizedSearchCV(text_clf, parameters, n_iter=n_iter_search)
    gs_clf = gs_clf.fit(X_train, y_train)
    text = gs_clf.predict([text])
    return text[0]


def url_to_text(url):
    urltext = pathlib.Path(url).suffix
    list_of_doc = ['.txt', '.PDF', '.pdf', '.csv', '.odp', '.ods', '.odt', '.rtf', '.docx', '.doc', '.xlsx', '.xls',
                   '.xlr', '.pptx', '.xml']
    if urltext not in list_of_doc:
        try:
            url = url
            response = requests.get(url, "lxml", headers={'User-Agent': 'Mozilla/5.0'})
            text = get_text(response.text)
            remove_n = text.replace('\n', ' ')
            remove_r = remove_n.replace('\r', ' ')
            remove_t = remove_r.replace('\t', ' ')
            clean_text = re.sub(' +', ' ', remove_t)
            remove_leading_space = clean_text.lstrip()
            return remove_leading_space
        except:
            return ' '
    else:
        if url.endswith('.txt'):
            try:
                data = urllib.request.urlopen(url)
                text = ''
                for line in data:
                    decoded_line = line.decode("utf-8")
                    text += decoded_line
                data = textclean(text)
                return data
            except:
                return ' '
        elif url.endswith('.PDF') or url.endswith('.pdf'):
            try:
                r = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                f = io.BytesIO(r.content)
                pdfFile = PyPDF2.PdfFileReader(f)
                No_of_pages = pdfFile.getNumPages()
                pages_content = ''
                for i in range(No_of_pages):
                    page = pdfFile.getPage(i)
                    page_content = page.extractText()
                    pages_content += page_content
                test = textclean(pages_content)
                return test
            except:
                return ' '
        elif url.endswith('.csv'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.txt')
                open(file_path, 'wb').write(myfile.content)
                source = open(file_path, "r")
                decoded_file = source.read()
                text = textclean(decoded_file)
                return text
            except:
                return ' '
        elif url.endswith('.odt'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.odt')
                open(file_path, 'wb').write(myfile.content)
                text = textract.process(file_path)
                esr = text.decode(encoding="utf-8")
                data = textclean(esr)
                return data
            except:
                return ' '
        elif url.endswith('.odp'):
            try:
                from odf import text, teletype
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.odp')
                open(file_path, 'wb').write(myfile.content)
                textdoc = load(file_path)
                allparas = textdoc.getElementsByType(text.P)
                linelenght = len(allparas)
                texts = ''
                for line in range(linelenght):
                    test = teletype.extractText(allparas[line])
                    texts += test
                data = textclean(texts)
                return data
            except:
                return ' '
        elif url.endswith('.ods'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.ods')
                open(file_path, 'wb').write(myfile.content)
                data1 = get_data(file_path)
                alphanumeric = [character for character in str(data1) if
                                character.isalnum() or character.isspace()]
                alphanumeric = "".join(alphanumeric)
                data = textclean(
                    alphanumeric)
                return data
            except:
                return ' '
        elif url.endswith('.docx'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.docx')
                open(file_path, 'wb').write(myfile.content)
                text = docx2txt.process(file_path)
                data = textclean(text)
                return data
            except:
                return ' '
        elif url.endswith('.doc'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                data = urllib.request.urlopen(url)
                with open('media/temp/temp.doc', 'wb') as output:
                    output.write(data.read())
                text = textract.process('media/temp/temp.doc')
                text = text.decode("utf-8")
                data = textclean(text)
                return data
            except:
                return ' '
        elif url.endswith('.xlr'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.xlr')
                xlrfile = pandas.read_excel('media/temp/temp.xlr')
                json_str = xlrfile.to_json(orient='records')
                data = textclean(json_str)
                return data
            except:
                return ' '
        elif url.endswith('.xlsx'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.xlsx')
                open(file_path, 'wb').write(myfile.content)
                df = pandas.read_excel(file_path)
                df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
                # data = self.textclean(text)
                return df
            except:
                return ' '
        elif url.endswith('.xls'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.xls')
                open(file_path, 'wb').write(myfile.content)
                df = pandas.read_excel('media/temp/temp.xls')
                df = df[~df.isin([np.nan, np.inf, -np.inf]).any(1)]
                os.remove("media/temp/temp.xls")
                # data = self.textclean(df)
                return df
            except:
                return ' '
        elif url.endswith('.xml'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                data = urllib.request.urlopen(url)
                with open('media/temp/temp.xml', 'wb') as output:
                    output.write(data.read())
                with open('media/temp/temp.xml', 'r') as f:
                    data = f.read()
                bsdata = BeautifulSoup(data, "xml")
                test = str(bsdata)
                cleanr = re.compile('<.*?>')
                cleantext = re.sub(cleanr, '', test)
                data = textclean(cleantext)
                return data
            except:
                return ' '
        elif url.endswith('.pptx'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.pptx')
                open(file_path, 'wb').write(myfile.content)
                prs = Presentation(file_path)
                text_runs = ''
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if not shape.has_text_frame:
                            continue
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text_runs += run.text
                data = textclean(text_runs)
                return data
            except:
                return ' '
        elif url.endswith('.rtf'):
            try:
                myfile = requests.get(url, headers={'User-Agent': 'Mozilla/5.0'})
                file_path = os.path.join(BASE_DIR, 'media/temp/temp.rtf')
                open(file_path, 'wb').write(myfile.content)
                text = open(file_path, 'r')
                rtf = text.read()
                txt = rtf_to_text(rtf)
                data = textclean(txt)
                return data
            except:
                return ' '


def textclean(text):
    remove_n = text.replace('\n', ' ')
    remove_r = remove_n.replace('\r', ' ')
    remove_t = remove_r.replace('\t', ' ')
    clean_text = re.sub(' +', ' ', remove_t)
    return remove_t


"""Background Jobs Functions"""

from background_task import background


@background(schedule=10)
def Documen_job(Keywords):
    print('Batch Added')
    try:
        is_success = Website_crawler(Keywords)
        print("success ful document job")
        os.system("pm2 restart Innovation_Db")
        if is_success == 'Error':
            os.system("pm2 restart Innovation_Db")
    except Exception as e:
        os.system("pm2 restart Innovation_Db")
        print(e)

from rest_framework import serializers
from searchBar.models import KeywordSearch

class KeywordSearchSerializers(serializers.ModelSerializer):
    search = serializers.CharField(source='keyWord')
    purpose = serializers.SerializerMethodField()
    technology = serializers.SerializerMethodField()
    date = serializers.DateTimeField(source='created_at')

    def get_purpose(self,obj):
        return []
    
    def get_technology(self,obj):
        return []

    class Meta:
        model = KeywordSearch
        fields = ('search','description','purpose','technology','status','date')